#!/usr/bin/env python3
"""Native (editable) PowerPoint comparing the tool's EM-Bayesian G-BVAR global solve with the
classical Pesaran GVAR — in the MFFM management-deck visual style (MadridBlue / Ink / Muted palette,
title bars, take-away boxes). Palette + slide helpers are reused from ../build_native_pptx.py.
Only the two TikZ diagrams are images; every bar / table cell / bullet is a real shape/text run.

Usage:  python3 build_compare_pptx.py [out.pptx] [figures_dir]
"""
import re, sys, os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__))
OUT = sys.argv[1] if len(sys.argv) > 1 else os.path.join(HERE, 'gbvar_vs_pesaran.pptx')
FIG = sys.argv[2] if len(sys.argv) > 2 else os.path.join(HERE, 'figures')

# ---- MFFM palette (identical to build_native_pptx.py) ----
MADRID=RGBColor(28,54,99); MADRID2=RGBColor(48,91,150); INK=RGBColor(35,35,35)
MUTED=RGBColor(105,105,105); SOFT=RGBColor(235,241,249); SOFT2=RGBColor(224,233,246)
MIDG=RGBColor(210,214,221); LIGHT=RGBColor(244,245,247)
RED=RGBColor(150,45,45); GREEN=RGBColor(39,116,75); WHITE=RGBColor(255,255,255)

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]
def slide(): return prs.slides.add_slide(BLANK)

# ---------------- inline-markup helpers (from the MFFM builder) ----------------
def clean(t):
    for a,b in [('---','—'),('--','–'),(r'\&','&'),(r'\%','%'),(r'\,',' '),('~',' '),
                ('``','“'),("''",'”'),(r'$\rightarrow$','→'),(r'$\cdot$','·'),
                (r'$\sim$','~'),(r'\textbullet','•'),('\\ ',' ')]:
        t=t.replace(a,b)
    return t
def parse(s):
    out=[]; i=0
    for m in re.finditer(r'\\(key|gap|emph|grn)\{([^}]*)\}', s):
        if m.start()>i: out.append((clean(s[i:m.start()]),'n'))
        k=m.group(1); inner=clean(m.group(2)); out.append((inner,k)); i=m.end()
    if i<len(s): out.append((clean(s[i:]),'n'))
    return out
def runs(para, segs, size, color=INK):
    para.line_spacing=1.03
    for text,k in segs:
        r=para.add_run(); r.text=text; f=r.font; f.name='Calibri'; f.size=Pt(size)
        if k=='key': f.bold=True; f.color.rgb=MADRID
        elif k=='gap': f.color.rgb=RED
        elif k=='grn': f.bold=True; f.color.rgb=GREEN
        elif k=='emph': f.italic=True; f.color.rgb=color
        else: f.color.rgb=color

def box(s, l,t,w,h, fill=None, line=None, lw=0.75, rounded=False):
    shp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
                           Inches(l),Inches(t),Inches(w),Inches(h))
    shp.shadow.inherit=False
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb=fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb=line; shp.line.width=Pt(lw)
    return shp
def tbox(s, l,t,w,h, anchor=MSO_ANCHOR.TOP):
    tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=Inches(0.05); tf.margin_right=Inches(0.05); tf.margin_top=Inches(0.02); tf.margin_bottom=Inches(0.02)
    return tf
def run(p,text,size,color=INK,bold=False,italic=False,font='Calibri'):
    r=p.add_run(); r.text=text; f=r.font; f.name=font; f.size=Pt(size); f.color.rgb=color; f.bold=bold; f.italic=italic; return r

def title_bar(s, text):
    box(s,0,0,13.333,0.82, fill=MADRID)
    tf=tbox(s,0.32,0,12.7,0.82, anchor=MSO_ANCHOR.MIDDLE)
    run(tf.paragraphs[0], clean(text), 19, WHITE, bold=True)

def takeaway(s, label, body_segs, y=6.42):
    b=box(s,0.5,y,12.333,0.86, fill=SOFT, line=MADRID, lw=0.9, rounded=True)
    tf=b.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    tf.margin_left=Inches(0.18); tf.margin_right=Inches(0.18)
    p=tf.paragraphs[0]; run(p, clean(label)+':  ', 14, MADRID, bold=True)
    runs(p, body_segs, 13, INK)
    return b

def header(s, l,t,w, text, color=MADRID, size=13, align=PP_ALIGN.LEFT):
    tf=tbox(s,l,t,w,0.35); p=tf.paragraphs[0]; p.alignment=align
    run(p, clean(text), size, color, bold=True)

def caption(s, l,t,w, segs, align=PP_ALIGN.CENTER, size=11):
    tf=tbox(s,l,t,w,0.4); p=tf.paragraphs[0]; p.alignment=align; runs(p, segs, size, MUTED)

def bullets(s, l,t,w,h, items, size=12.5, isp=6):
    tf=tbox(s,l,t,w,h)
    for i,(marker,segs) in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.space_after=Pt(isp)
        lead={'plus':('+ ',GREEN,True),'dash':('– ',RED,True),'dot':('▸ ',MADRID,True),
              '1':('1. ',MADRID,True),'2':('2. ',MADRID,True)}.get(marker)
        if lead:
            r=p.add_run(); r.text=lead[0]; f=r.font; f.name='Calibri'; f.size=Pt(size); f.bold=lead[2]; f.color.rgb=lead[1]
        runs(p, segs, size, INK)

def fit_img(s, name, bl,bt,bw,bh, border=True):
    """Place figures/<name>.png centred inside the (bl,bt,bw,bh) box, preserving aspect."""
    path=os.path.join(FIG,name+'.png')
    iw,ih=Image.open(path).size; ar=iw/ih
    w=bw; h=w/ar
    if h>bh: h=bh; w=h*ar
    l=bl+(bw-w)/2; t=bt+(bh-h)/2
    s.shapes.add_picture(path, Inches(l),Inches(t),Inches(w),Inches(h))
    if border: box(s,l,t,w,h, line=MIDG, lw=0.75, rounded=True)
    return l,t,w,h

# ============================================================ SLIDES
# ---------------------------------------------------------------- 1 TITLE
s=slide()
tf=tbox(s,1.0,1.55,11.3,1.9); p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
run(p,'The Global Solve, Two Ways', 34, MADRID, bold=True)
q=tf.add_paragraph(); q.alignment=PP_ALIGN.CENTER; q.space_before=Pt(6)
run(q,'EM-Bayesian G-BVAR  vs.  the classical Pesaran GVAR', 21, MADRID2, bold=True)
box(s,5.42,3.35,2.5,0.03, fill=MADRID)
tf=tbox(s,1.5,3.6,10.3,1.0); p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
runs(p,[('Same skeleton — country VARX* models linked by trade-weighted foreign variables — two ways to estimate it and to solve the coupled global system.','n')],15,INK)
tf=tbox(s,1.5,5.05,10.3,1.7); p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
run(p,'Mátyás Farkas',14,INK)
for txt,sz,col,it in [('International Monetary Fund',12,MUTED,False),
                      ('Technical background note — MFFM global solve',11.5,MUTED,False),
                      ('Views expressed are my own and do not necessarily represent the IMF.',9,MUTED,True)]:
    r=tf.add_paragraph(); r.alignment=PP_ALIGN.CENTER; r.space_before=Pt(7)
    run(r,txt,sz,col,italic=it)

# ---------------------------------------------------------------- 2 SUMMARY
s=slide(); title_bar(s,'Executive summary — both are GVARs; the tool generalizes the classical one')
bullets(s,0.7,1.15,12.0,5.0,[
 ('dot',parse(r'\key{Both are GVARs.} Every economy is a VAR in its own variables plus \key{trade-weighted foreign “stars”} (partners’ activity, prices, rates); the country models are stacked and solved into \key{one global system}.')),
 ('dot',parse(r'\key{(a) Bayesian, not OLS.} Each bloc is a GLP hierarchical-Minnesota BVAR (shrunk toward a random walk, tightness λ / sum-of-coefficients μ sampled) — not the classical per-bloc OLS VARX* / reduced-rank VECMX*.')),
 ('dot',parse(r'\key{(b) It completes the data.} An \key{EM loop} with a Durbin–Koopman simulation smoother fills ragged edges and COVID-as-missing macros (financials kept observed) — the classical GVAR estimates each bloc on its own span, then couples on a common window (truncate / dummy, \gap{no imputation}).')),
 ('dot',parse(r'\key{(c) The US is endogenous.} It carries its own ROW*_US stars (n 37→40) so the \key{world feeds back} onto it; the classical GVAR treats the US as the dominant \emph{source} of the weakly-exogenous global commons.')),
 ('dot',parse(r'\key{(d) An affine fixed point.} The coupled solve reduces to x = Mx + d ⇒ x = (I−M)⁻¹d with a \key{pre-shipped inverse}, coupling by projection step; the classical solve stacks and inverts the contemporaneous matrix G₀ once.')),
],size=13.5, isp=9)
takeaway(s,'In one line',parse(r'A \key{superset, not a rewrite}: the same VARX* + trade-weighted-star GVAR skeleton, made Bayesian, gap-tolerant, US-endogenous, and solved as an affine deviation fixed point.'))

# ---------------------------------------------------------------- 3 PESARAN DIAGRAM
s=slide(); title_bar(s,'What the classical GVAR does (Pesaran–Schuermann–Weiner; Dees–di Mauro–Pesaran–Smith)')
tf=tbox(s,0.6,0.92,12.1,0.4); p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
runs(p,parse(r'Estimate each country \key{separately}, build weakly-exogenous stars from observed partners, then \key{stack and invert G₀ once}.'),12.5)
fit_img(s,'fig_pesaran',0.5,1.38,12.333,4.9)
takeaway(s,'The classical recipe',parse(r'\key{N separate} VARX* / VECMX* blocs, each on its own sample · fixed-weight, weakly-exogenous stars · a \gap{common solve window} (ragged edges truncated, no imputation) · US a source · one G₀ inversion.'))

# ---------------------------------------------------------------- 4 G-BVAR DIAGRAM
s=slide(); title_bar(s,'How the EM-Bayesian G-BVAR differs — completion loop, Bayesian blocs, endogenous US')
tf=tbox(s,0.6,0.92,12.1,0.4); p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
runs(p,parse(r'An \key{EM + DK-smoother} loop completes the panels and rebuilds the stars; GLP-Bayesian blocs; the US is \key{endogenous}; the coupled solve is an \key{affine fixed point} with a pre-computed inverse.'),12)
fit_img(s,'fig_gbvar',0.5,1.38,12.333,4.9)
takeaway(s,'The four generalizations',parse(r'\grn{+} Bayesian GLP shrinkage · \grn{+} EM/DK data completion · \grn{+} endogenous US (world feeds back) · \grn{+} affine (I−M)⁻¹d solve.'))

# ---------------------------------------------------------------- 5 COMPARISON TABLE
s=slide(); title_bar(s,'Side by side — the same GVAR skeleton, five points of departure')
COL0_L,COL0_W=0.5,2.35
COL1_L,COL1_W=2.90,4.98
COL2_L,COL2_W=7.93,4.90
HY,HH=1.05,0.52
def hcell(l,w,text,fill):
    box(s,l,HY,w,HH,fill=fill,line=WHITE,lw=1.0)
    tf=tbox(s,l,HY,w,HH,anchor=MSO_ANCHOR.MIDDLE); p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    run(p,text,13,WHITE,bold=True)
hcell(COL0_L,COL0_W,'',MADRID)
hcell(COL1_L,COL1_W,'Classical Pesaran GVAR',MADRID2)
hcell(COL2_L,COL2_W,'EM-Bayesian G-BVAR (this tool)',MADRID)
# cell content (kept as data so it stays grounded + editable)
cells={
 'Estimation':(
   parse(r'Each country \key{separately}: per-bloc OLS VARX* — or reduced-rank VECMX* ML. Fixed weights, no shrinkage prior.'),
   parse(r'\key{GLP hierarchical-Minnesota BVAR} per bloc; shrunk to a random walk; λ, μ sampled, not fixed; posterior-mean companion exported.')),
 'Foreign variables':(
   parse(r'Stars  x*_i = Σ w_ij x_j  over partners j≠i; from \key{observed} data; \key{weakly exogenous} (tested per bloc).'),
   parse(r'3 stars per bloc (DEM / INF / FIN); \key{per-quarter renormalised} weights; estimated as ordinary columns, \key{pinned affinely} at solve.')),
 'Missing data':(
   parse(r'Each bloc on its \key{own span}; the \key{coupled solve} uses a common window (ragged edges truncated); COVID dummied — \gap{no model-based imputation}.'),
   parse(r'\key{EM + DK simulation smoother} completes ragged & COVID-missing macros; financials kept observed; coupling by \key{projection step}.')),
 'Dominant unit (US)':(
   parse(r'US = \key{dominant unit} — source of the weakly-exogenous global commons. Its foreign \key{financial} stars are restricted (x*_US restricted), so the financial channel is \gap{one-directional — the world does not feed back through finance}.'),
   parse(r'US \key{endogenous}: own ROW*_US stars, n 37→40; \key{world feeds back}; US drives the 4 global commons.')),
 'Global solve':(
   parse(r'Stack  G₀ x_t = a + Σ Gℓ x_(t−ℓ) + u;  \key{invert G₀ once} → reduced-form GVAR (GIRFs).'),
   parse(r'Affine deviation  x = Mx + d ⇒ x = (I−M)⁻¹d;  \key{pre-shipped inverse}; unique since ρ(M) = 0.947.')),
}
order=['Estimation','Foreign variables','Missing data','Dominant unit (US)','Global solve']
ry=HY+HH; rh=(6.28-ry)/len(order)
for i,name in enumerate(order):
    y=ry+i*rh; band=LIGHT if i%2==0 else WHITE
    box(s,COL0_L,y,COL0_W,rh,fill=SOFT2,line=WHITE,lw=1.0)
    tf=tbox(s,COL0_L+0.05,y,COL0_W-0.1,rh,anchor=MSO_ANCHOR.MIDDLE); p=tf.paragraphs[0]
    run(p,name,12,MADRID,bold=True)
    lseg,rseg=cells[name]
    box(s,COL1_L,y,COL1_W,rh,fill=band,line=WHITE,lw=1.0)
    tf=tbox(s,COL1_L+0.02,y,COL1_W-0.04,rh,anchor=MSO_ANCHOR.MIDDLE)
    tf.margin_left=Inches(0.12); tf.margin_right=Inches(0.12); runs(tf.paragraphs[0],lseg,11)
    box(s,COL2_L,y,COL2_W,rh,fill=band,line=WHITE,lw=1.0)
    tf=tbox(s,COL2_L+0.02,y,COL2_W-0.04,rh,anchor=MSO_ANCHOR.MIDDLE)
    tf.margin_left=Inches(0.12); tf.margin_right=Inches(0.12); runs(tf.paragraphs[0],rseg,11)
takeaway(s,'Read the columns',parse(r'Left = the textbook GVAR; right = the \key{same construction} made Bayesian, gap-tolerant, US-endogenous and affine.'),y=6.45)

# ---------------------------------------------------------------- 6 NUMBERS + WHAT IT BUYS
s=slide(); title_bar(s,'The shipped numbers, and what the generalization buys')
header(s,0.6,1.05,6.0,'The shipped 33-bloc coupled system')
bullets(s,0.6,1.5,6.3,4.6,[
 ('dot',parse(r'\key{ρ(M) = 0.946626 < 1} → the fixed point is unique (frozen-anchor legacy: ρ ≈ 0.70).')),
 ('dot',parse(r'Augmented state \key{dim = 2060} = 33·3·20 + 4·20 (stars stacked with US-global deltas) — verified against the shipped artifact.')),
 ('dot',parse(r'\key{50 economies} = 33 GVAR-coupled blocs + 17 euro members; lags p = 3, horizon H = 20.')),
 ('dot',parse(r'EM loop: MAXIT ≤ 12, TOL = 0.03; GLP Gibbs nDraws = 1000, nBurn = 500.')),
 ('dot',parse(r'\key{Pre-shipping the M + (I−M)⁻¹ artifact} skips the ≈19-min in-browser smoother-probe M-build (24 blocs; more at 33) and the ≈46-min in-JS inverse — offline the inverse itself is ≈1 s (BLAS). Warm in-browser solve ≈4.5 s; a US-pinned structure falls back to the inverse-free Neumann path (fresh M).')),
],size=12.5, isp=9)
header(s,7.15,1.05,5.7,'What it buys')
bullets(s,7.15,1.5,5.7,4.6,[
 ('plus',parse(r'\key{Current, shrunk estimates} on ragged, COVID-scarred data — no common-window truncation.')),
 ('plus',parse(r'A genuine \key{world→US feedback} channel the classical GVAR excludes by construction.')),
 ('plus',parse(r'\key{One affine solve}, computed once from a pre-built inverse, reused for every scenario.')),
 ('plus',parse(r'\key{Coupling by projection step}, not calendar date — a bloc ending a quarter early is \emph{completed}, not dropped everywhere.')),
],size=12.5, isp=9)
box(s,7.0,1.45,0.012,4.55, fill=MIDG)  # thin divider
takeaway(s,'Bottom line',parse(r'The classical GVAR is the \key{fixed-weight, common-window, exogenous-US special case}; the tool relaxes all three and keeps the solve linear.'))

prs.save(OUT); print('saved',OUT,'-',len(prs.slides._sldIdLst),'slides')
