#!/usr/bin/env python3
"""Standalone 2-slide EDITABLE pptx: the government debt accounting equations exactly as
implemented (computeFiscalBlock), in the MFFM management-deck style. All equations are native
text runs (Cambria Math) with TRUE OOXML sub/superscripts — no images, everything editable.

Usage: python3 build_debt_eqn_slides.py [out.pptx]
"""
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUT = sys.argv[1] if len(sys.argv) > 1 else 'debt_accounting_equations.pptx'
MADRID = RGBColor(28, 54, 99); INK = RGBColor(35, 35, 35); MUTED = RGBColor(105, 105, 105)
SOFT = RGBColor(235, 241, 249); MIDG = RGBColor(210, 214, 221); LIGHT = RGBColor(244, 245, 247)
WHITE = RGBColor(255, 255, 255)

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)

def box(s, l, t, w, h, fill=None, line=None, lw=0.75, rounded=True):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
                             Inches(l), Inches(t), Inches(w), Inches(h))
    shp.shadow.inherit = False
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(lw)
    return shp

def tbox(s, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h)); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.06); tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    return tf

def run(p, text, size, color=INK, bold=False, italic=False, font='Calibri'):
    r = p.add_run(); r.text = text; f = r.font
    f.name = font; f.size = Pt(size); f.color.rgb = color; f.bold = bold; f.italic = italic
    return r

def mrun(p, text, size=16, italic=False, base=None, color=INK, bold=False):
    r = p.add_run(); r.text = text; f = r.font
    f.name = 'Cambria Math'; f.size = Pt(size); f.color.rgb = color; f.italic = italic; f.bold = bold
    if base == 'sub': r._r.get_or_add_rPr().set('baseline', '-25000')
    elif base == 'sup': r._r.get_or_add_rPr().set('baseline', '30000')
    return r

def eqn(p, parts, size=16):
    # parts: (text, italic, base) triples; base in {None,'sub','sup'}
    for t, it, base in parts: mrun(p, t, size, italic=it, base=base)

def title_bar(s, text):
    box(s, 0, 0, 13.333, 0.82, fill=MADRID, rounded=False)
    tf = tbox(s, 0.32, 0, 12.7, 0.82, anchor=MSO_ANCHOR.MIDDLE)
    run(tf.paragraphs[0], text, 19, WHITE, bold=True)

def card(s, y, h, label, parts, size=16, accent=False, note=None):
    box(s, 0.6, y, 12.13, h, fill=SOFT if accent else LIGHT,
        line=MADRID if accent else MIDG, lw=1.4 if accent else 0.75)
    run(tbox(s, 0.78, y + 0.07, 11.8, 0.3).paragraphs[0], label, 11.5, MADRID, bold=True)
    tf = tbox(s, 0.6, y + 0.34, 12.13, h - 0.4 - (0.3 if note else 0), anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    eqn(p, parts, size=size)
    if note:
        tf2 = tbox(s, 0.78, y + h - 0.34, 11.6, 0.3)
        run(tf2.paragraphs[0], note, 9.5, MUTED)

# convenience fragments
def v(t): return (t, True, None)          # italic math variable
def n(t): return (t, False, None)         # upright (numbers, operators, parens)
def sub(t): return (t, True, 'sub')       # italic subscript
def subn(t): return (t, False, 'sub')     # upright subscript
def sup(t): return (t, False, 'sup')      # upright superscript

# ============ SLIDE 1: one-quarter accounting ============
s = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(s, 'Government debt accounting — one quarter, as implemented')

card(s, 1.00, 1.02, 'Nominal GDP growth (quarterly, decimal; ℓ = 100·ln of the level path)',
     [v('g'), sub('t'), sup('n'), n('  =  [ ('), v('ℓ'), sub('t'), sup('gdp'), n(' − '),
      v('ℓ'), sub('t−1'), sup('gdp'), n(') + ('), v('ℓ'), sub('t'), sup('defl'), n(' − '),
      v('ℓ'), sub('t−1'), sup('defl'), n(') ] / 100')], size=16)

card(s, 2.14, 1.02, 'Flow updates (annual ratios, % of GDP)',
     [v('pb'), sub('t'), n('  =  '), v('rev'), sub('t'), n(' − '), v('pexp'), sub('t'),
      n('      '), v('int'), sub('t'), n('  =  '), v('i'), sub('t'), sup('eff'), n(' · '),
      v('d'), sub('t−1'), n(' / 100'),
      n('      '), v('def'), sub('t'), n('  =  '), v('int'), sub('t'), n(' − '), v('pb'), sub('t')], size=16)

card(s, 3.28, 1.30, 'The debt update (the roll-forward; seeded at the last observed ratio d₀)',
     [v('d'), sub('t'), n('  =  '), v('d'), sub('t−1'), n(' · '),
      n('(1 + '), v('i'), sub('t'), sup('eff'), n('/400)'), n(' / '), n('(1 + '), v('g'), sub('t'), sup('n'), n(')'),
      n('  −  '), v('pb'), sub('t'), n('/4'),
      n('  +  '), v('sfa'), n('/4')], size=19, accent=True,
     note='sfa = the structural stock-flow constant (long-run mean; sfaDefault: US 1.75, EA 0.33, JP 2.23, UK 3.15, CA 3.19, CN 4.02 % of GDP) — deliberately not the forecast SFA column.')

card(s, 4.72, 1.30, 'Change per quarter — the snowball decomposition',
     [n('Δ'), v('d'), sub('t'), n('  =  [ (1 + '), v('i'), sub('t'), sup('eff'), n('/400)/(1 + '),
      v('g'), sub('t'), sup('n'), n(') − 1 ] · '), v('d'), sub('t−1'),
      n('  −  '), v('pb'), sub('t'), n('/4  +  '), v('sfa'), n('/4'),
      n('   ≈   '), n('('), v('i'), sub('t'), sup('eff'), n('/400 − '), v('g'), sub('t'), sup('n'),
      n(')/(1 + '), v('g'), sub('t'), sup('n'), n(') · '), v('d'), sub('t−1'),
      n('  −  '), v('pb'), sub('t'), n('/4  +  '), v('sfa'), n('/4')], size=14.5,
     note='Debt rises mechanically whenever the quarterized effective rate exceeds quarterly nominal growth (i − g > 0).')

tf = tbox(s, 0.6, 6.30, 12.13, 0.85)
p = tf.paragraphs[0]
run(p, 'Estimated drivers:  ', 11, MADRID, bold=True)
run(p, 'rev, pexp (% of GDP), the effective rate iᵉᶠᶠ (% p.a.), SFA — ordinary BVAR variables. ', 11, INK)
run(p, 'Derived by identity:  ', 11, MADRID, bold=True)
run(p, 'pb, interest, deficit, Δd, and the debt ratio d — never estimated freely: the identity has a near-unit root (1+i)/(1+g), so free estimation is explosive.', 11, INK)

# ============ SLIDE 2: evolution over time ============
s = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(s, 'Debt evolution over time — compounding, and the exact inversion')

card(s, 1.05, 1.00, 'The one-quarter compounding factor',
     [v('R'), sub('t'), n('  ≡  (1 + '), v('i'), sub('t'), sup('eff'), n('/400)'),
      n(' / (1 + '), v('g'), sub('t'), sup('n'), n(')')], size=18)

card(s, 2.20, 1.72, 'The T-quarter evolution (the recursion unrolled)',
     [v('d'), subn('T'), n('  =  '), v('d'), subn('0'), n(' · '),
      n('∏'), subn('t=1..T'), n(' '), v('R'), sub('t'),
      n('   +   '), n('∑'), subn('t=1..T'), n(' ( '), n('∏'), subn('s=t+1..T'), n(' '), v('R'), sub('s'), n(' )'),
      n(' · ( '), v('sfa'), n(' − '), v('pb'), sub('t'), n(' ) / 4')], size=17, accent=True,
     note='The inherited stock compounds at the interest–growth ratio; each quarter’s net flow is carried forward at all subsequent R. With R ≈ 1 (near-unit root) small persistent flows accumulate one-for-one — and the compounding is convex in the drivers, so the per-draw debt distribution is right-skewed.')

card(s, 4.07, 1.30, 'Exact inversion — the primary balance that delivers a target debt path',
     [v('pb'), sub('t'), n('  =  4 · [ ('), v('d'), sub('t−1'), sup('tgt'),
      n(' + '), v('int'), sub('t'), n('/4) / (1 + '), v('g'), sub('t'), sup('n'), n(')'),
      n('  −  '), v('d'), sub('t'), sup('tgt'), n(' ]  +  '), v('sfa')], size=17,
     note='This is what a dragged Debt-to-GDP path conditions on: the required pb is pinned as the linear combination rev − pexp in the smoother.')

tf = tbox(s, 0.6, 5.62, 12.13, 1.5)
p = tf.paragraphs[0]
run(p, 'Reading the dynamics:  ', 11.5, MADRID, bold=True)
run(p, 'three forces move the ratio — the snowball (i − g, automatic), the primary balance (policy), and the stock-flow constant (measurement/perimeter). '
       'A long-run growth anchor changes the terminal gⁿ and therefore owns the debt slope through the snowball — arithmetic, not a fiscal reform. '
       'Uncertainty bands push every posterior draw through this same identity, so the debt band widens exactly as fast as budget-and-growth uncertainty implies.', 11.5, INK)

prs.save(OUT)
print(f'saved {OUT} - {len(prs.slides._sldIdLst)} slides')
