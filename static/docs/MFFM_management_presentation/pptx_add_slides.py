#!/usr/bin/env python3
"""Add the debt-accumulation equation to the fiscal slide + a GVAR 'how it's solved' background slide.
Native shapes/text with TRUE sub/superscripts. Run on a freshly-generated base pptx."""
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

PPTX=sys.argv[1]
MADRID=RGBColor(28,54,99); MADRID2=RGBColor(48,91,150); INK=RGBColor(35,35,35)
MUTED=RGBColor(105,105,105); SOFT=RGBColor(235,241,249); MIDG=RGBColor(210,214,221)
RED=RGBColor(150,45,45); WHITE=RGBColor(255,255,255); LIGHT=RGBColor(244,245,247)
prs=Presentation(PPTX)

def box(s,l,t,w,h,fill=None,line=None,lw=0.75,rounded=False):
    shp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    shp.shadow.inherit=False
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb=fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb=line; shp.line.width=Pt(lw)
    return shp
def tbox(s,l,t,w,h,anchor=MSO_ANCHOR.TOP):
    tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=Inches(0.06); tf.margin_right=Inches(0.06); tf.margin_top=Inches(0.02); tf.margin_bottom=Inches(0.02)
    return tf
def run(p,text,size,color=INK,bold=False,italic=False,font='Calibri'):
    r=p.add_run(); r.text=text; f=r.font; f.name=font; f.size=Pt(size); f.color.rgb=color; f.bold=bold; f.italic=italic; return r
def mrun(p,text,size=16,italic=False,base=None,color=INK):
    r=p.add_run(); r.text=text; f=r.font; f.name='Cambria Math'; f.size=Pt(size); f.color.rgb=color; f.italic=italic
    if base=='sub': r._r.get_or_add_rPr().set('baseline','-25000')
    elif base=='sup': r._r.get_or_add_rPr().set('baseline','30000')
    return r
def eqn(p, parts, size=16):
    for t,it,base in parts: mrun(p,t,size,italic=it,base=base)
def title_bar(s,text):
    box(s,0,0,13.333,0.82,fill=MADRID); tf=tbox(s,0.32,0,12.7,0.82,anchor=MSO_ANCHOR.MIDDLE)
    run(tf.paragraphs[0],text,19,WHITE,bold=True)
def takeaway(s,label,body,y=6.38):
    b=box(s,0.5,y,12.333,0.92,fill=SOFT,line=MADRID,lw=0.9,rounded=True); tf=b.text_frame
    tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE; tf.margin_left=Inches(0.18); tf.margin_right=Inches(0.18)
    p=tf.paragraphs[0]; run(p,label+':  ',14,MADRID,bold=True)
    for text,st in body: run(p,text,13.5, MADRID if st=='key' else (RED if st=='gap' else INK), bold=(st=='key'), italic=(st=='emph'))

# ---------- 1) DEBT EQUATION on the fiscal slide (index 5) ----------
s=prs.slides[5]; ey=3.5; ew=5.5
box(s,0.6,ey,ew,1.5,fill=LIGHT,line=MIDG,lw=0.75,rounded=True)
run(tbox(s,0.78,ey+0.1,ew-0.3,0.32).paragraphs[0],'The debt dynamics, one line:',11.5,MADRID,bold=True)
tf=tbox(s,0.6,ey+0.46,ew-0.02,0.55,anchor=MSO_ANCHOR.MIDDLE); p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
eqn(p,[('d',True,None),('t',False,'sub'),('  =  ',False,None),('d',True,None),('t−1',False,'sub'),
       (' · ',False,None),('(1+',False,None),('i',True,None),('t',False,'sub'),(') / (1+',False,None),('g',True,None),('t',False,'sub'),(')',False,None),
       ('  −  ',False,None),('pb',True,None),('t',False,'sub'),('  +  ',False,None),('sfa',True,None),('t',False,'sub')],size=17)
tf=tbox(s,0.7,ey+1.04,ew-0.2,0.4); p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
run(p,'d = debt/GDP · i = eff. interest rate · g = nominal GDP growth · pb = primary balance · sfa = stock-flow adj.',8.7,MUTED)

# ---------- 2) GVAR BACKGROUND SLIDE (appended) ----------
s=prs.slides.add_slide(prs.slide_layouts[6])
title_bar(s,'Background — how the coupled global solve works (GVAR)')
tf=tbox(s,0.6,1.15,6.5,4.6)
pts=[('The idea',True,MADRID,13),
     ('Each economy is estimated with its own variables PLUS a set of trade-weighted “foreign” variables — its partners’ activity, prices and rates.',False,INK,12.5),
     ('Stacking all 50 economies, those foreign links tie every model to every other into one big system.',False,INK,12.5),
     ('We then solve for the one state where every economy’s forecast is mutually consistent with its partners’ — the fixed point.',False,INK,12.5),
     ('Computed once from a pre-built inverse (fast), or iterated to convergence — either way, no country is solved in isolation.',False,MUTED,12)]
for i,(t,b,c,sz) in enumerate(pts):
    p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.space_after=Pt(8)
    if i>0: run(p,'▸  ',12,MADRID,bold=True)
    run(p,t,sz,c,bold=b)
rx=7.5; rw=5.3
def mathbox(y,label,parts,h=1.0):
    box(s,rx,y,rw,h,fill=LIGHT,line=MIDG,lw=0.75,rounded=True)
    run(tbox(s,rx+0.18,y+0.07,rw-0.36,0.3).paragraphs[0],label,10.5,MADRID,bold=True)
    tf=tbox(s,rx,y+0.36,rw,h-0.42,anchor=MSO_ANCHOR.MIDDLE); p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; eqn(p,parts,size=17)
def darrow(y):
    a=s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,Inches(rx+rw/2-0.08),Inches(y),Inches(0.16),Inches(0.26)); a.shadow.inherit=False; a.fill.solid(); a.fill.fore_color.rgb=MADRID2; a.line.fill.background()
mathbox(1.35,'Foreign (trade-weighted) variables',
        [('x',True,None),('i',True,'sub'),('*',False,'sup'),('  =  ',False,None),('Σ',False,None),('j',False,'sub'),(' ',False,None),('w',True,None),('ij',True,'sub'),(' ',False,None),('x',True,None),('j',True,'sub')])
darrow(2.42)
mathbox(2.72,'Stack every economy + its links',[('x',True,None),('  =  ',False,None),('a',True,None),('  +  ',False,None),('M',True,None),(' ',False,None),('x',True,None)])
darrow(3.79)
mathbox(4.09,'Solve the fixed point (mutually consistent world)',[('x',True,None),('  =  (I − ',False,None),('M',True,None),(')',False,None),('−1',False,'sup'),(' ',False,None),('a',True,None)])
run(tbox(s,rx,5.28,rw,0.5,anchor=MSO_ANCHOR.TOP).paragraphs[0]  ,'M holds the trade-weighted cross-country links; a is each economy’s own dynamics.',9.5,MUTED)
takeaway(s,'In plain terms',[('No country moves alone — the solve finds the world path where every economy’s reaction to its partners is ',''),('mutually consistent','key'),('.','')])

prs.save(PPTX); print('saved',PPTX,'-',len(prs.slides._sldIdLst),'slides')
