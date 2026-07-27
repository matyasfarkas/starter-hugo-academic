#!/usr/bin/env python3
"""Rebuild the MFFM management deck as a NATIVE (editable) PowerPoint: title bars, roadmap strips,
text, bullets and take-away boxes are real shapes/text; only the chart screenshots are images."""
import re, sys, os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

FIG = sys.argv[2]  # figures dir
OUT = sys.argv[1]

MADRID=RGBColor(28,54,99); MADRID2=RGBColor(48,91,150); INK=RGBColor(35,35,35)
MUTED=RGBColor(105,105,105); SOFT=RGBColor(235,241,249); MIDG=RGBColor(210,214,221)
RED=RGBColor(150,45,45); GREEN=RGBColor(39,116,75); WHITE=RGBColor(255,255,255)

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]
def slide(): return prs.slides.add_slide(BLANK)

def clean(t):
    for a,b in [('---','—'),('--','–'),(r'\&','&'),(r'\%','%'),(r'\,',' '),('~',' '),
                ('``','“'),("''",'”'),(r'$\rightarrow$','→'),(r'$\cdot$','·'),
                (r'$\sim$','~'),(r'$-0.8$','−0.8'),(r'\textbullet','•'),('\\ ',' ')]:
        t=t.replace(a,b)
    return t
def parse(s):
    out=[]; i=0
    for m in re.finditer(r'\\(key|gap|emph|todo)\{([^}]*)\}', s):
        if m.start()>i: out.append((clean(s[i:m.start()]),'n'))
        k=m.group(1); inner=clean(m.group(2))
        if k=='todo': inner='['+inner+']'
        out.append((inner,k)); i=m.end()
    if i<len(s): out.append((clean(s[i:]),'n'))
    return out
def runs(para, segs, size, color=INK):
    para.line_spacing=1.02
    for text,k in segs:
        r=para.add_run(); r.text=text; f=r.font; f.name='Calibri'; f.size=Pt(size)
        if k=='key': f.bold=True; f.color.rgb=MADRID
        elif k=='gap': f.color.rgb=RED
        elif k=='emph': f.italic=True; f.color.rgb=color
        elif k=='todo': f.italic=True; f.color.rgb=RED; f.size=Pt(size-1)
        else: f.color.rgb=color

def box(s, l,t,w,h, fill=None, line=None, lw=0.75, rounded=False):
    shp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
                           Inches(l),Inches(t),Inches(w),Inches(h))
    shp.shadow.inherit=False
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb=fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb=line; shp.line.width=Pt(lw)
    return shp

def tbox(s, l,t,w,h, anchor=MSO_ANCHOR.TOP):
    tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=Inches(0.05); tf.margin_right=Inches(0.05); tf.margin_top=Inches(0.02); tf.margin_bottom=Inches(0.02)
    return tf

def title_bar(s, text):
    box(s,0,0,13.333,0.82, fill=MADRID)
    tf=tbox(s,0.32,0,12.7,0.82, anchor=MSO_ANCHOR.MIDDLE)
    p=tf.paragraphs[0]; r=p.add_run(); r.text=clean(text); f=r.font
    f.name='Calibri'; f.size=Pt(19); f.bold=True; f.color.rgb=WHITE

ROAD=['Baseline','Conditioning','Global solve','Extensions','Conclusion']
def roadmap(s, active):
    bw,bh,gap=1.55,0.34,0.30; n=5; total=n*bw+(n-1)*gap; x0=(13.333-total)/2; y=0.96
    for i,lab in enumerate(ROAD):
        x=x0+i*(bw+gap); on=(i+1)==active
        b=box(s,x,y,bw,bh, fill=(MADRID if on else SOFT), line=MADRID, lw=0.6, rounded=True)
        tf=b.text_frame; tf.word_wrap=False; p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
        r=p.add_run(); r.text=lab; f=r.font; f.name='Calibri'; f.size=Pt(9)
        f.color.rgb=(WHITE if on else MADRID); f.bold=on
        if i<n-1:
            ar=s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x+bw+0.02),Inches(y+bh/2-0.045),Inches(gap-0.04),Inches(0.09))
            ar.shadow.inherit=False; ar.fill.solid(); ar.fill.fore_color.rgb=MADRID2; ar.line.fill.background()

def takeaway(s, label, body_segs, y=6.38):
    b=box(s,0.5,y,12.333,0.92, fill=SOFT, line=MADRID, lw=0.9, rounded=True)
    tf=b.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    tf.margin_left=Inches(0.18); tf.margin_right=Inches(0.18)
    p=tf.paragraphs[0]; r=p.add_run(); r.text=clean(label)+':  '; f=r.font
    f.name='Calibri'; f.bold=True; f.size=Pt(14); f.color.rgb=MADRID
    runs(p, body_segs, 13.5, INK)
    return b

def img_card(s, name, l,t,w):
    path=os.path.join(FIG,name+'.png')
    iw,ih=Image.open(path).size; h=w*ih/iw
    s.shapes.add_picture(path, Inches(l),Inches(t),Inches(w),Inches(h))
    box(s,l,t,w,h, line=MIDG, lw=0.75, rounded=True)  # card outline
    return h

def header(s, l,t,w, text, color=MADRID, size=13, align=PP_ALIGN.CENTER):
    tf=tbox(s,l,t,w,0.35); p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=clean(text); f=r.font; f.name='Calibri'; f.bold=True; f.size=Pt(size); f.color.rgb=color

def caption(s, l,t,w, segs, align=PP_ALIGN.CENTER, size=10.5):
    tf=tbox(s,l,t,w,0.4); p=tf.paragraphs[0]; p.alignment=align; runs(p, segs, size, MUTED)

def para_block(s, l,t,w,h, paras, size=12.5, anchor=MSO_ANCHOR.TOP):
    # paras: list of segment-lists
    tf=tbox(s,l,t,w,h, anchor=anchor)
    for i,segs in enumerate(paras):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.space_after=Pt(8); runs(p, segs, size, INK)

def bullets(s, l,t,w,h, items, size=12.5):
    # items: list of (marker, segment-list); marker in {'plus','dash','1','2','dot',None}
    tf=tbox(s,l,t,w,h)
    for i,(marker,segs) in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.space_after=Pt(5)
        if marker=='plus': lead=('+ ',GREEN,True)
        elif marker=='dash': lead=('– ',RED,True)
        elif marker=='1': lead=('1. ',MADRID,True)
        elif marker=='2': lead=('2. ',MADRID,True)
        elif marker=='dot': lead=('▸ ',MADRID,True)
        else: lead=None
        if lead:
            r=p.add_run(); r.text=lead[0]; f=r.font; f.name='Calibri'; f.size=Pt(size); f.bold=lead[2]; f.color.rgb=lead[1]
        runs(p, segs, size, INK)

# ============================ SLIDES ============================
# 1 TITLE
s=slide()
tf=tbox(s,1.0,1.4,11.3,1.6); p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
r=p.add_run(); r.text='The Monetary–Financial–Fiscal Macromodel (MFFM)'; f=r.font; f.name='Calibri'; f.bold=True; f.size=Pt(34); f.color.rgb=MADRID
box(s,5.4,2.7,2.5,0.03, fill=MADRID)
tf=tbox(s,1.5,2.95,10.3,1.0); p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
runs(p,[('A global projection and scenario tool — one coherent “what-if” for 50 economies, and its policy implications.','n')],18,INK)
tf=tbox(s,1.5,4.5,10.3,1.6); p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
r=p.add_run(); r.text='Mátyás Farkas'; r.font.size=Pt(15); r.font.color.rgb=INK; r.font.name='Calibri'
for txt,sz,col in [('International Monetary Fund',12,MUTED),('Senior Management briefing   •   [date / venue]',12,MUTED),
                   ('Views expressed are my own and do not necessarily represent the IMF.',9,MUTED)]:
    q=tf.add_paragraph(); q.alignment=PP_ALIGN.CENTER; q.space_before=Pt(8)
    r=q.add_run(); r.text=txt; r.font.size=Pt(sz); r.font.color.rgb=col; r.font.name='Calibri'
    if 'date' in txt: r.font.italic=True; r.font.color.rgb=RED

# 2 MOTIVATION
s=slide(); title_bar(s,"Why now: the next generation of the Fund's global macro-financial model")
header(s,0.6,1.05,5.8,'Today: the GMF (Vitek)',MADRID,13,PP_ALIGN.LEFT)
bullets(s,0.6,1.5,5.9,3.6,[
 (None,[('Served the Fund well as a global macro-financial workhorse.','n')]),
 ('dash',parse(r'\gap{Not updated in estimation} --- \todo{MF/Jesper: specifics}')),
 ('dash',parse(r'\gap{Missing FSAP / Article IV features} --- \todo{MF/Jesper: list the needs}')),
 ('dash',parse(r'\todo{MF/Jesper: any further gaps}')),
])
header(s,6.9,1.05,5.9,'What surveillance now needs',MADRID,13,PP_ALIGN.LEFT)
bullets(s,6.9,1.5,5.9,3.6,[
 ('plus',[('Current estimates, refreshed each cycle.','n')]),
 ('plus',[('Fast, transparent scenarios for FSAP & Article IV.','n')]),
 ('plus',[('Genuine cross-border transmission — shocks travel and feed back across economies.','n')]),
 ('plus',parse(r'\todo{MF/Jesper: the specific FSAP/AIV features}')),
])
takeaway(s,'The case',parse(r'\todo{MF/Jesper: one-line motivation --- e.g. “The MFFM supersedes the GMF: re-estimated, scenario-ready, and built for FSAP--Article IV work.”}'))

# 3 ADVANTAGES
s=slide(); title_bar(s,'What the MFFM adds — five things it does that a static model cannot')
bullets(s,0.7,1.25,12.0,4.9,[
 ('plus',parse(r'\key{Re-estimated and current.} Bayesian VAR baselines re-estimated on the latest data, for \key{50 economies} --- one consistent methodology everywhere.')),
 ('plus',parse(r'\key{Interactive scenarios.} Build a “what-if” by hand, or from an outside forecast, in minutes --- the model keeps every variable mutually consistent.')),
 ('plus',parse(r'\key{Real cross-border amplification.} A coupled \emph{global} solve --- a shock in one economy travels through trade and finance to the rest, not a sum of separate country runs.')),
 ('plus',parse(r'\key{Uses what we already produce.} WEO, consensus and desk forecasts enter as \emph{ingredients}, at a trust level you choose.')),
 ('plus',parse(r'\key{Honest and reproducible.} Uncertainty bands throughout; a whole session --- every country, every assumption --- saved and shared as a single file.')),
],size=14)
takeaway(s,'In one line',parse(r'The same trusted models, made \key{current, interactive, global, and shareable}.'))

# 4 OVERVIEW
s=slide(); title_bar(s,'One tool: a coherent global scenario, its spillovers, and the policy trade-offs')
para_block(s,0.7,0.95,12.0,0.6,[parse(r'A point-and-click platform for \key{50 economies} --- the Fund’s workhorse models (Bayesian VAR baselines, a cross-country GVAR, DSGE policy) behind a single interface. Today’s tour, in five steps:')],size=12.5)
flow=[('Baseline','data-driven forecast, 50 economies'),('Conditioning','drag a path, or import a forecast'),
      ('Global solve','one country moves the rest'),('Extensions','policy, anchors, external forecasts'),
      ('Conclusion','bands · export · share · decide')]
fw,fgap,fh=2.28,0.28,1.0; ftot=5*fw+4*fgap; fx0=(13.333-ftot)/2; fy=1.75
for i,(h,d) in enumerate(flow):
    x=fx0+i*(fw+fgap); b=box(s,x,fy,fw,fh, fill=SOFT, line=MADRID, lw=0.7, rounded=True)
    tf=b.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; r=p.add_run(); r.text=h; r.font.bold=True; r.font.size=Pt(11); r.font.color.rgb=MADRID; r.font.name='Calibri'
    q=tf.add_paragraph(); q.alignment=PP_ALIGN.CENTER; r=q.add_run(); r.text=d; r.font.size=Pt(9); r.font.color.rgb=INK; r.font.name='Calibri'
    if i<4:
        ar=s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x+fw+0.01),Inches(fy+fh/2-0.05),Inches(fgap-0.02),Inches(0.1))
        ar.shadow.inherit=False; ar.fill.solid(); ar.fill.fore_color.rgb=MADRID2; ar.line.fill.background()
header(s,0.7,3.1,6.0,'Why it matters',MADRID,12.5,PP_ALIGN.LEFT)
bullets(s,0.7,3.5,6.3,2.8,[
 ('plus',parse(r'\key{Speed} --- a coherent global “what-if” in an afternoon, not a forecasting round.')),
 ('plus',parse(r'\key{Coherence} --- every number is mutually consistent, at home \emph{and} across borders.')),
 ('plus',parse(r'\key{Reuse} --- folds in the forecasts we already produce (WEO, SPF, desk).')),
])
header(s,7.1,3.1,5.9,'What you can ask it',MADRID,12.5,PP_ALIGN.LEFT)
bullets(s,7.1,3.5,5.9,2.8,[
 ('dot',[('If US growth surprises — who else moves, and by how much?','n')]),
 ('dot',[('If the Fed holds rates higher for longer, what happens here?','n')]),
 ('dot',[('If we take the WEO as given, what does it imply for everything else?','n')]),
])

# helper for a text-left / figure-right section slide
def sec_lr(num, active, title, left_paras, figname, figw, cap_segs, take_label, take_segs, left_w=0.44, todo_line=None, left_head=None, left_bul=None):
    s=slide(); title_bar(s, '①②③④⑤'[num-1]+'  '+title); roadmap(s,active)
    lw=left_w*12.4; rx=0.6+lw+0.5; rw=13.333-rx-0.5
    if left_head: header(s,0.6,1.55,lw+0.2,left_head,MADRID,12.5,PP_ALIGN.LEFT); ly=1.95
    else: ly=1.7
    if left_bul: bullets(s,0.6,ly,lw+0.2,3.6,left_bul,size=13)
    else: para_block(s,0.6,ly,lw+0.2,3.6,left_paras,size=13)
    fw_in=figw*rw; fx=rx+(rw-fw_in)/2
    h=img_card(s,figname,fx,1.75,fw_in)
    caption(s,rx,1.75+h+0.08,rw,cap_segs)
    takeaway(s,take_label,take_segs)
    if todo_line:
        tf=tbox(s,0.5,7.32,12.3,0.3); p=tf.paragraphs[0]; runs(p,parse(todo_line),9.5)

# 5  (1a) BASELINE — centered figure
s=slide(); title_bar(s,'①  Start with the baseline — no inputs, the model’s own read of every economy'); roadmap(s,1)
tf=tbox(s,1.0,1.5,11.3,0.5); p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
runs(p,parse(r'\key{No assumptions required.} Each economy is a Bayesian VAR, linked into one global system.'),13)
h=img_card(s,'s1_macro_triad',1.9,2.05,9.5)
caption(s,1.9,2.05+h+0.06,9.5,parse(r'US inflation, policy rate and GDP --- one consistent forecast, with uncertainty bands.'))
takeaway(s,'The baseline',parse(r'A ready-made, internally consistent forecast for all \key{50 economies} --- the common starting point every scenario departs from.'))

# 6 (1b) DEBT
sec_lr(1,1,'Fiscal accounting, built in — the debt path drops out of the same run',
 [parse(r'Revenue, spending, the interest bill and the deficit are projected \key{together} --- and debt-to-GDP follows from the \key{budget identity}.'),
  parse(r'No separate spreadsheet: the fiscal story is \key{consistent with the macro forecast by construction}.')],
 's1_debt_gdp',0.9, parse(r'US debt-to-GDP rises to ~143\% by 2031, with band.'),
 'Why it matters', parse(r'Debt sustainability drops out of the \emph{same} run --- the macro and fiscal stories can never contradict each other.'))

# 7 (2a) ENERGY two-up — leads (clean stagflation = best first impression)
s=slide(); title_bar(s,'②  Tell it what you believe — the model works out the rest, consistently'); roadmap(s,2)
header(s,0.6,1.6,6.0,'You pin',MADRID,14)
h=img_card(s,'fig_energy_oil',1.4,2.05,4.7); caption(s,0.6,2.05+h+0.06,6.0,parse(r'An oil-price path --- the size and shape of the spike.'))
header(s,6.7,1.6,6.0,'Out comes',MADRID,14)
h2=img_card(s,'fig_energy_macro',7.5,2.05,4.7); caption(s,6.7,2.05+h2+0.06,6.0,parse(r'US stagflation: \key{prices up, output down}.'))
takeaway(s,'The power move',parse(r'State a couple of things you believe --- the model works out a \key{full, model-consistent response} across every variable, no re-estimation.'))

# 8 (2b) CREDIT — second example
sec_lr(2,2,'Another shock, the same move — a credit crunch',
 None,'fig_credit_macro',0.92, parse(r'US GDP response to the two credit pins, with bands.'),
 'Fit the pins to the story', parse(r'Two pins for a crunch, a few more for a richer scenario --- either way the model keeps everything else \emph{consistent}.'),
 left_w=0.42, left_head='Recipe · Credit crunch',
 left_bul=[('1',parse(r'Pin \key{spreads jump}.')),('2',parse(r'Pin \key{lending falls}.')),
           (None,parse(r'Two pins in; the model works out a \key{full, coherent response} across every variable.'))])

# 9 (3) CLIMAX two-up
s=slide(); title_bar(s,'③  A shock abroad, a hit at home — the part a country-only read cannot see'); roadmap(s,3)
header(s,0.6,1.6,6.0,'Country-only read',MUTED,14)
h=img_card(s,'fig_partner_energy_dev_direct',1.55,2.05,4.4); caption(s,0.6,2.05+h+0.06,6.0,parse(r'Euro-area GDP vs baseline: \key{flat}.'))
header(s,6.7,1.6,6.0,'Global solve',MADRID,14)
h2=img_card(s,'fig_partner_energy_dev_coupled',7.65,2.05,4.4); caption(s,6.7,2.05+h2+0.06,6.0,parse(r'The same shock arrives: \key{−0.8pp} by 2031.'))
takeaway(s,'Why the tool earns its keep',parse(r'Read one economy at a time and this hit is \key{invisible}; only the coupled solve makes the cross-border spillover \key{show up} — and sizes it.'))

# 10 (4a) ANCHORS
sec_lr(4,4,'Pin the destination — and the whole path tilts toward it',
 [parse(r'Know only \key{where things should end up} --- the inflation target, potential growth, the neutral rate?'),
  parse(r'Set that \key{long-run anchor}; the model tilts the path toward it, gently, \key{without breaking the near term}.')],
 'fig_ext_anchor',0.9, parse(r'US inflation: the tail lands on the anchor (dashed); near term untouched.'),
 'Long-run anchors', parse(r'Set where a variable \key{settles}; the path re-shapes to reach it, near term intact.'))

# 11 (4b) WEO
sec_lr(4,4,'Fold in an outside forecast as an ingredient — at a trust level you choose',
 [parse(r'A WEO number, a consensus figure, a desk view? \key{Import it directly} --- as one more ingredient, not a hard override.'),
  parse(r'You set the \key{trust dial}: from “just a hint” to “pin it exactly”. Everything else stays \key{consistent} around it.')],
 'fig_ext_weo',0.86, parse(r'US GDP: teal diamonds are the imported WEO path.'),
 'External forecasts', parse(r'Reuse the WEO and desk forecasts we \emph{already} produce --- as inputs, at a \key{trust level you control}.'))

# 12 (4c) MMB
sec_lr(4,4,'Ask “what if the central bank does X?” — a genuine policy counterfactual',
 [parse(r'This tells you \key{what a policy choice does} --- a \key{structural} answer from a fully specified DSGE model.'),
  parse(r'Layer a \key{policy shock} on any scenario; the model traces the \key{counterfactual path}.')],
 'fig_ext_mmb',0.8, parse(r'US policy rate: counterfactual vs. baseline scenario.'),
 'Policy counterfactuals', parse(r'From “what will happen” to “\key{what should we do}” --- a structural (MMB/DSGE) policy experiment.'),
 todo_line=r'\todo{on the roadmap: deeper MMB library, cross-country WEO propagation, external-balance/NFA satellite}')

# 13 (5) CONCLUSION
s=slide(); title_bar(s,'⑤  One tool now gives us a current, global, and shareable read on policy — we ask you to adopt it'); roadmap(s,5)
tf=tbox(s,0.6,1.6,6.3,0.35); r=tf.paragraphs[0].add_run(); r.text='What we built, in three lines:'; r.font.bold=True; r.font.size=Pt(13); r.font.color.rgb=INK; r.font.name='Calibri'
bullets(s,0.6,2.0,6.3,3.5,[
 ('plus',parse(r'\key{Current & consistent} --- every economy sits on the same up-to-date vintage, so the numbers add up across desks')),
 ('plus',parse(r'\key{Genuinely global} --- shocks travel across 50 economies through trade and finance, not one country at a time')),
 ('plus',parse(r'\key{Our forecasts, built in} --- WEO and staff views enter as dials, and any run is one-click \key{shareable and reproducible}')),
],size=12.5)
para_block(s,7.1,1.6,5.9,4.4,[
 parse(r'\key{The arc:} a live baseline → \key{condition} it on the story you care about → let it \key{propagate globally} → layer in \key{policy and external-balance} extensions.'),
 parse(r'\key{The ask:} \todo{MF/Jesper: endorse the tool, resource the last-mile hardening, and adopt it into the FSAP / Article IV workflow.}'),
 parse(r'It is \key{live today}, covering \key{50 economies} — ready to pilot on the next surveillance cycle.'),
],size=12.5)
takeaway(s,'The one line',parse(r'One live, shareable tool turns scattered forecasts into a single consistent, global answer to “what if?”.'))

prs.save(OUT); print('saved',OUT,'-',len(prs.slides._sldIdLst),'slides')
